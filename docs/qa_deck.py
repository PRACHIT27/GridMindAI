"""Geometry QA for the generated deck.

No LibreOffice on this machine, so instead of rendering and eyeballing, this
reads the shape geometry back out of the .pptx and checks the three defects that
actually bite in a generated deck:

  1. anything outside the slide, or inside the 0.5" margin
  2. text that needs more lines than its box has room for
  3. text boxes that overlap each other

Width estimation is per-font (monospace is exact; the proportional fonts use
measured average advance widths for their character mix), so a flag here means
the text really is close to the edge.
"""
from __future__ import annotations

import math
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

DECK = Path(__file__).resolve().parent / "GridMind_Submission_Deck.pptx"
W, H, MARGIN = 13.333, 7.5, 0.5

# average advance width as a fraction of point size, for mixed-case prose
ADVANCE = {"Courier New": 0.600, "Arial": 0.512, "Calibri": 0.468}


def inches(v) -> float:
    return Emu(int(v)).inches


def wrapped_lines(text: str, box_w: float, size: float, font: str, bold: bool) -> int:
    """Greedy word wrap at the estimated pixel width."""
    adv = ADVANCE.get(font, 0.50) * (1.06 if bold else 1.0) * size / 72.0
    total = 0
    for para in text.split("\n"):
        if not para.strip():
            total += 1
            continue
        line, n = 0.0, 1
        for word in para.split(" "):
            w = (len(word) + 1) * adv
            if line + w > box_w and line > 0:
                n += 1
                line = w
            else:
                line += w
        total += n
    return total


def walk(shape):
    if shape.shape_type == 6:  # group
        for sub in shape.shapes:
            yield from walk(sub)
    else:
        yield shape


def main() -> int:
    prs = Presentation(DECK)
    problems: list[str] = []

    for idx, slide in enumerate(prs.slides, 1):
        boxes = []
        for shape in slide.shapes:
            for sh in walk(shape):
                x, y = inches(sh.left), inches(sh.top)
                w, h = inches(sh.width), inches(sh.height)

                # ---- 1. bounds
                if x < -0.01 or y < -0.01 or x + w > W + 0.01 or y + h > H + 0.01:
                    problems.append(
                        f"slide {idx}: off-slide  x={x:.2f} y={y:.2f} "
                        f"w={w:.2f} h={h:.2f}")
                elif x < MARGIN - 0.01 or y < MARGIN - 0.01 or \
                        x + w > W - MARGIN + 0.01 or y + h > H - MARGIN + 0.01:
                    problems.append(
                        f"slide {idx}: inside 0.5\" margin  x={x:.2f} y={y:.2f} "
                        f"right={x + w:.2f} bottom={y + h:.2f}")

                if not sh.has_text_frame:
                    continue
                txt = sh.text_frame.text
                if not txt.strip():
                    continue

                # ---- 2. text fit
                runs = [r for p in sh.text_frame.paragraphs for r in p.runs]
                if not runs:
                    continue
                f = runs[0].font
                size = f.size.pt if f.size else 18
                font = f.name or "Calibri"
                spacing = next((p.line_spacing.pt for p in sh.text_frame.paragraphs
                                if p.line_spacing is not None), None)
                lh = (spacing or size * 1.21) / 72.0
                lines = wrapped_lines(txt, w, size, font, bool(f.bold))
                need = lines * lh
                if need > h + 0.06:
                    problems.append(
                        f"slide {idx}: text overflow ({lines} lines x {lh * 72:.0f}pt "
                        f"= {need:.2f}\" in a {h:.2f}\" box)  \"{txt[:58]}...\"")
                if y + need > H - 0.2:
                    problems.append(
                        f"slide {idx}: text runs off the bottom (ends {y + need:.2f}\")"
                        f"  \"{txt[:48]}...\"")
                boxes.append((x, y, w, min(need, h), txt))

        # ---- 3. text-on-text overlap
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                ax, ay, aw, ah, at = boxes[i]
                bx, by, bw, bh, bt = boxes[j]
                ox = min(ax + aw, bx + bw) - max(ax, bx)
                oy = min(ay + ah, by + bh) - max(ay, by)
                if ox > 0.05 and oy > 0.05:
                    problems.append(
                        f"slide {idx}: text overlap {ox:.2f}x{oy:.2f}\"  "
                        f"\"{at[:26]}\" / \"{bt[:26]}\"")

    if problems:
        print(f"{len(problems)} issue(s):\n")
        for p in problems:
            print("  " + p)
        return 1
    print(f"clean — {len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
          f"bounds / text fit / overlap all pass")
    return 0


if __name__ == "__main__":
    sys.exit(main())
