"""GridMind — hackathon submission deck.

Order: what the product does (2-6), then how it is built (7-11), then coverage
and limits (12-14). The rubric weights still drive how much room each part gets,
but a judge should understand the product before seeing a service diagram.

    python docs/make_deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "GridMind_Submission_Deck.pptx"

C = dict(
    bg=RGBColor(0x0E, 0x13, 0x19), panel=RGBColor(0x18, 0x20, 0x2C),
    panel2=RGBColor(0x1E, 0x27, 0x35), edge=RGBColor(0x2A, 0x35, 0x43),
    ink=RGBColor(0xEE, 0xF2, 0xF7), dim=RGBColor(0x9A, 0xA8, 0xBA),
    faint=RGBColor(0x6B, 0x7A, 0x8D), cy=RGBColor(0x4D, 0xD6, 0xFF),
    grn=RGBColor(0x34, 0xD3, 0x99), amb=RGBColor(0xFB, 0xBF, 0x24),
    red=RGBColor(0xF8, 0x71, 0x71), vio=RGBColor(0xA7, 0x8B, 0xFA),
)
H, B, M = "Arial", "Calibri", "Courier New"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["bg"]
    return s


def text(s, x, y, w, h, body, *, size=13, color="dim", font=B, bold=False,
         italic=False, align=PP_ALIGN.LEFT, spacing=None, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = Pt(spacing)
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = C[color]
    return tb


def card(s, x, y, w, h, fill="panel", line="edge"):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = C[fill]
    sh.line.color.rgb = C[line]
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def dot(s, x, y, d, color):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = C[color]
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def title(s, head, sub=None):
    text(s, 0.6, 0.50, 12.1, 0.64, head, size=29, bold=True, color="ink", font=H)
    if sub:
        text(s, 0.6, 1.16, 12.1, 0.36, sub, size=14, color="dim")


def bullets(s, x, y, w, items, gap=0.42, size=12.5, color="dim", bullet="grn"):
    for i, it in enumerate(items):
        yy = y + i * gap
        dot(s, x, yy + 0.09, 0.11, bullet)
        text(s, x + 0.28, yy, w - 0.28, gap, it, size=size, color=color)


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


# =============================================================== 1. title
s = slide()
text(s, 0.75, 1.90, 8.4, 1.15, "GridMind", size=60, bold=True, color="ink", font=H)
text(s, 0.78, 3.06, 8.6, 0.5, "Capacity planning for high-density data centers",
     size=21, color="cy")
text(s, 0.78, 3.74, 8.3, 1.2,
     "Every new workload has to clear power, cooling, floor space and budget before it can be "
     "placed. Today four teams check those four things separately, over several days. GridMind "
     "checks them together and returns a plan the site can actually build.", size=14, spacing=21)
card(s, 9.55, 1.90, 3.1, 3.55)
for i, (k, v, col, big) in enumerate([
        ("Track", "The Fortified\nEnterprise Fleet", "ink", True),
        ("Live", "gridmind-wuvfpvopoq\n-uk.a.run.app", "cy", False),
        ("Built on", "Gemini 3.5 on Vertex AI\nCloud Run, Firestore", "dim", False)]):
    y = 2.18 + i * 1.10
    text(s, 9.85, y, 2.6, 0.26, k, size=10.5, color="faint")
    text(s, 9.85, y + 0.25, 2.6, 0.68, v, size=13 if big else 11.5, bold=big,
         color=col, font=H if big else B, spacing=16)
notes(s, "One line on what it is, then move. The card on the right lets judges see the stack "
         "without me listing it.")

# ======================================================= 2. what it does
s = slide()
title(s, "What GridMind does", "A capacity request goes in. A plan the site can build comes out.")
steps = [
    ("Submit a request", "Rack count, power draw per rack, cooling type, weight, budget "
                         "ceiling. A form, not free text.", "cy"),
    ("Four assessors run in parallel", "Power, cooling, facilities and cost. Each one reads "
                                       "only its own data and answers one question.", "grn"),
    ("The orchestrator reconciles", "It checks whether the four answers describe one buildable "
                                    "plan. When they do not, it sends each assessor the others' "
                                    "positions and asks again.", "amb"),
    ("You get a decision and a report", "Target zone, the plan to get there, and what every "
                                        "assessor said beside the numbers it was shown.", "vio"),
]
for i, (n, d, col) in enumerate(steps):
    y = 1.75 + i * 1.18
    card(s, 0.6, y, 12.1, 1.02)
    dot(s, 0.92, y + 0.30, 0.42, col)
    text(s, 0.92, y + 0.36, 0.42, 0.32, str(i + 1), size=16, bold=True, color="bg",
         font=H, align=PP_ALIGN.CENTER)
    text(s, 1.58, y + 0.20, 3.5, 0.32, n, size=15, bold=True, color="ink")
    text(s, 5.25, y + 0.22, 7.2, 0.62, d, size=12.5, spacing=17)
text(s, 0.6, 6.55, 12.1, 0.36,
     "About ninety seconds end to end. The same four checks take a site several days today.",
     size=13, color="cy")
notes(s, "Walk the four steps. This is the whole product in one slide.")

# ============================================= 3. inputs and outputs
s = slide()
title(s, "What you submit, and what you get back")
card(s, 0.6, 1.72, 5.9, 4.6)
text(s, 0.95, 1.98, 5.2, 0.3, "The request", size=15, bold=True, color="cy")
bullets(s, 0.95, 2.44, 5.2, [
    "Workload profile, or a custom rack spec",
    "Number of racks",
    "Power draw per rack, in kW",
    "Cooling requirement, air or liquid",
    "Weight per rack, in kg",
    "Budget ceiling",
], gap=0.52, bullet="cy")
text(s, 0.95, 5.62, 5.2, 0.6,
     "Every field is a number or a fixed choice. Nothing an operator types is passed to a model "
     "as instructions.", size=11.5, color="faint", spacing=16)
card(s, 6.8, 1.72, 5.9, 4.6)
text(s, 7.15, 1.98, 5.2, 0.3, "The result", size=15, bold=True, color="grn")
bullets(s, 7.15, 2.44, 5.2, [
    "A target zone, or a refusal with the reason",
    "The plan: which racks, what work, how long",
    "Each assessor's verdict and its evidence",
    "A decision report as PDF",
    "The full decision history, kept and searchable",
], gap=0.52, bullet="grn")
text(s, 7.15, 5.20, 5.2, 1.0,
     "The report is the part operators asked about first. It shows what each assessor was given "
     "next to what it concluded, so an approval can be checked rather than trusted.",
     size=11.5, color="faint", spacing=16)
notes(s, "The point of the right-hand column is that the output is auditable, not just an "
         "answer.")

# ================================================== 4. why it exists
s = slide()
title(s, "Why it exists",
      "Four teams, four inboxes, and nobody who checks the four answers together")
for i, (name, what, col) in enumerate([
        ("Power engineering", "breaker capacity, spare circuits, grid events", "cy"),
        ("Cooling and thermal ops", "heat removal, efficiency, water permits", "grn"),
        ("Facilities and DCOps", "rack space, floor loading, install crews", "amb"),
        ("Finance and procurement", "budget, cost per GPU-hour, capex", "vio")]):
    y = 1.75 + i * 0.84
    card(s, 0.6, y, 6.5, 0.72)
    dot(s, 0.88, y + 0.22, 0.28, col)
    text(s, 1.32, y + 0.10, 2.9, 0.28, name, size=13.5, bold=True, color="ink")
    text(s, 1.32, y + 0.38, 5.3, 0.28, what, size=11.5)
card(s, 7.5, 1.75, 5.2, 3.50, "panel2")
text(s, 7.85, 2.02, 4.5, 0.85, "Each one is right.\nNobody owns the joint check.",
     size=19, bold=True, color="ink", font=H, spacing=25)
text(s, 7.85, 3.00, 4.5, 2.20,
     "Requests move from team to team by email and spreadsheet. Every approval is correct on its "
     "own axis, and no one asks whether the four approvals describe the same plan.\n\n"
     "What that leaves behind has a name in the industry: stranded capacity. Megawatts a site has "
     "paid for and cannot use.", size=12.5, spacing=19)
for i, (n, l, col) in enumerate([("days", "for four sequential approvals", "amb"),
                                 ("0", "teams who check them together", "red")]):
    text(s, 0.6 + i * 3.3, 5.45, 3.0, 0.7, n, size=38, bold=True, color=col, font=H)
    text(s, 0.6 + i * 3.3, 6.14, 3.0, 0.3, l, size=12)
text(s, 7.5, 5.45, 5.2, 0.7, "~90 sec", size=38, bold=True, color="grn", font=H)
text(s, 7.5, 6.14, 5.2, 0.4, "for the same four checks, run together", size=12)
notes(s, "Short slide. The next one shows the failure rather than describing it.")

# ================================================= 5. a real run
s = slide()
title(s, "A real run, round one",
      "Six racks of GB200 NVL72. 792 kW, liquid cooled, 1,360 kg per rack.")
for i, (who, zone, why, col) in enumerate([
        ("Power", "zone-a", "most electrical headroom", "red"),
        ("Cooling", "zone-c", "only zone that can remove the heat", "grn"),
        ("Facilities", "zone-b", "only zone with 6 free liquid racks", "amb"),
        ("Cost", "zone-b", "avoids a $48k per rack retrofit", "amb")]):
    x = 0.6 + i * 3.06
    card(s, x, 1.78, 2.92, 1.9)
    text(s, x + 0.26, 1.96, 2.4, 0.3, who, size=13, bold=True, color="dim")
    text(s, x + 0.26, 2.30, 2.4, 0.5, zone, size=25, bold=True, color=col, font=M)
    text(s, x + 0.26, 2.88, 2.42, 0.7, why, size=11, color="faint", spacing=15)
card(s, 0.6, 4.00, 12.1, 1.3, "panel2")
text(s, 0.95, 4.20, 11.4, 0.42,
     "Four correct answers. Three different rooms. No plan.",
     size=18, bold=True, color="ink", font=H)
text(s, 0.95, 4.68, 11.4, 0.42,
     "A system that counted votes would see four approvals and place the order. Then the racks "
     "arrive and the room cannot cool them.", size=13)
text(s, 0.6, 5.62, 12.1, 0.42,
     "This is the failure we built for. Not agents contradicting each other, but agents all "
     "being right and still adding up to nothing.", size=13.5, color="cy")
notes(s, "Say the four zones out loud, then pause before the next slide.")

# ============================================ 6. how the plan is found
s = slide()
title(s, "How the plan gets found",
      "Every assessor also reports the zones its own axis rules out")
for i, (who, zones) in enumerate([("Cooling rules out", "zone-a   zone-b   zone-d"),
                                  ("Power rules out", "zone-b   zone-d"),
                                  ("Facilities rules out", "zone-a   zone-d"),
                                  ("Cost rules out", "zone-a   zone-d")]):
    y = 1.88 + i * 0.62
    text(s, 0.6, y, 2.6, 0.36, who, size=13, align=PP_ALIGN.RIGHT)
    text(s, 3.35, y, 3.4, 0.36, zones, size=13, color="red", font=M)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.35), Inches(4.42), Inches(3.2), Pt(1.2))
ln.fill.solid()
ln.fill.fore_color.rgb = C["edge"]
ln.line.fill.background()
ln.shadow.inherit = False
text(s, 0.6, 4.60, 2.6, 0.36, "survives all four", size=13, bold=True, color="ink",
     align=PP_ALIGN.RIGHT)
text(s, 3.35, 4.54, 3.4, 0.5, "zone-c", size=20, bold=True, color="grn", font=M)
text(s, 0.6, 5.40, 6.1, 0.9,
     "The orchestrator does not count votes. It intersects the four exclusion lists and asks "
     "whether anything is left.", size=12.5, spacing=18)
card(s, 7.2, 1.78, 5.5, 4.55, "panel2")
text(s, 7.55, 2.04, 4.8, 0.8, "In round one, not one assessor had picked it.",
     size=19, bold=True, color="grn", font=H, spacing=25)
text(s, 7.55, 2.96, 4.8, 3.25,
     "Power wanted zone-a. Facilities and cost wanted zone-b. The zone that actually works was "
     "nobody's first choice.\n\n"
     "Asked again with the conflict attached, facilities volunteered something it had no reason "
     "to mention before: zone-c has five liquid-ready racks and more that can be plumbed in about "
     "fourteen hours each. Cost re-priced one retrofit at $48,000 instead of the $288,000 it had "
     "assumed for six.\n\n"
     "Final plan: zone-c, five ready racks plus one retrofit, fourteen hours of delay.",
     size=12, spacing=18)
notes(s, "This is the answer no single team can reach from inside its own silo. That is the "
         "reason the system is worth building.")

# ============================================== 7. system architecture
s = slide()
title(s, "System architecture", "Seven Cloud Run services, seven identities, five databases")


def lane(y, h, label, col):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.6), Inches(y), Inches(12.1), Inches(h))
    sh.adjustments[0] = 0.04
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x14, 0x1B, 0x25)
    sh.line.color.rgb = C[col]
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    text(s, 8.6, y + 0.08, 4.0, 0.26, label, size=10, bold=True, color=col,
         align=PP_ALIGN.RIGHT)


lane(1.68, 0.98, "PUBLIC", "faint")
card(s, 0.85, 1.88, 3.1, 0.62)
text(s, 1.05, 1.95, 2.8, 0.26, "gridmind, web tier", size=12.5, bold=True, color="cy")
text(s, 1.05, 2.20, 2.8, 0.26, "read only, no model access", size=10, color="faint")
text(s, 4.25, 2.00, 6.2, 0.45,
     "The only service reachable from the internet, and the least privileged one in the system.",
     size=11.5)
lane(2.82, 2.42, "PRIVATE VPC, ingress=internal, 404 from the internet", "vio")
for i, (n, d, col) in enumerate([("Orchestrator", "shared store only", "cy"),
                                 ("Gateway", "identity, routing, Model Armor", "amb"),
                                 ("Four assessors", "one store each", "grn")]):
    x = 0.85 + i * 3.95
    card(s, x, 3.10, 3.7, 0.78)
    text(s, x + 0.2, 3.19, 3.3, 0.28, n, size=13, bold=True, color=col)
    text(s, x + 0.2, 3.47, 3.3, 0.28, d, size=10.5, color="faint")
text(s, 0.85, 4.06, 11.6, 0.32,
     "A harness wraps every assessor: scoped read, schema-checked output, guardrail re-check, "
     "retry with correction, fail safe.", size=11.5, bold=True, color="amb")
text(s, 0.85, 4.42, 11.6, 0.6,
     "Gemini 3.5 on Vertex AI. flash-lite for each assessor, flash with a 4,096-token thinking "
     "budget for reconciliation.", size=11.5, color="vio")
lane(5.36, 0.96, "DATA", "grn")
for i, d in enumerate(["power-db", "cooling-db", "facilities-db", "cost-db", "shared-db"]):
    x = 0.85 + i * 2.36
    card(s, x, 5.54, 2.16, 0.6)
    text(s, x, 5.62, 2.16, 0.26, d, size=11, bold=True, color="grn", font=M,
         align=PP_ALIGN.CENTER)
    text(s, x, 5.86, 2.16, 0.24, "all others 403", size=9.5, color="red",
         align=PP_ALIGN.CENTER)
text(s, 0.6, 6.55, 12.1, 0.32,
     "Full diagram with the Google Cloud icons, plus the setup scripts, are in the repo.",
     size=12)
notes(s, "Do not read the boxes. The shape is the point: the most exposed service holds the "
         "least privilege, and everything else is off the internet.")

# ================================================= 8. data isolation
s = slide()
title(s, "Data isolation, enforced by the platform",
      "Not a convention in our code. An IAM condition Firestore itself applies.")
text(s, 0.6, 1.70, 12.1, 0.6,
     "Firestore Security Rules are not evaluated for server-side Admin SDK access. A design that "
     "scoped agents with per-collection rules would have enforced nothing at all, while looking "
     "correct in a diagram.", size=13.5, color="amb", spacing=20)
sec = [("Identity", "Its own store", "The other stores", "From the internet"),
       ("power, cooling, facilities, cost", "200", "403", "404"),
       ("orchestrator", "shared only", "403, all four", "404"),
       ("public web tier", "shared, read only", "403", "200, the only door")]
tbl = s.shapes.add_table(4, 4, Inches(0.6), Inches(2.5), Inches(12.1), Inches(2.1)).table
for w, col in zip((4.4, 2.5, 2.7, 2.5), range(4)):
    tbl.columns[col].width = Inches(w)
for r, row in enumerate(sec):
    tbl.rows[r].height = Inches(0.5)
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["panel2"] if r == 0 else C["panel"]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.14)
        run = cell.text_frame.paragraphs[0].add_run()
        run.text = val
        run.font.size = Pt(13)
        run.font.name = B
        run.font.bold = (r == 0) or (c > 0 and val[:3] in ("200", "403", "404"))
        run.font.color.rgb = (C["faint"] if r == 0 else C["ink"] if c == 0 else
                              C["red"] if val[:3] in ("403", "404") else
                              C["cy"] if "only door" in val else C["grn"])
card(s, 0.6, 4.95, 5.9, 1.55, "panel2")
text(s, 0.95, 5.15, 5.2, 0.3, "The orchestrator cannot read facility data",
     size=14.5, bold=True, color="ink", font=H)
text(s, 0.95, 5.50, 5.2, 0.9,
     "So \"it only sees the assessors' verdicts\" is a property of the platform, not a promise in "
     "a README. It is bound to one store and cannot open the others.", size=12, spacing=17)
card(s, 6.8, 4.95, 5.9, 1.55, "panel2")
text(s, 7.15, 5.15, 5.2, 0.3, "Twenty-one checks, run live in the demo",
     size=14.5, bold=True, color="ink", font=H)
text(s, 7.15, 5.50, 5.2, 0.6,
     "bash infra/iam/03_verify_isolation.sh\nbash scripts/demo_denial.sh",
     size=11.5, color="grn", font=M, spacing=17)
text(s, 7.15, 6.08, 5.2, 0.3, "Every denial comes from Google Cloud, not from our code.",
     size=11, color="faint")
notes(s, "The Security Rules finding is the strongest technical point in the project. Say it "
         "plainly. Fourteen isolation checks plus seven denial checks.")

# ========================================== 9. when an agent gets it wrong
s = slide()
title(s, "When an assessor gets it wrong",
      "The rubric asks how the system recovers if a worker agent loops or returns "
      "a hallucination")
card(s, 0.6, 1.90, 5.9, 2.2, "panel2")
text(s, 0.95, 2.08, 5.2, 0.28, "A verdict the model actually produced", size=12, color="faint")
text(s, 0.95, 2.40, 5.2, 1.1,
     '"status": "feasible",\n"target_zone": "zone-b",\n"reasoning": "Zone B has ample\n'
     '  headroom for this deployment."', size=11.5, color="ink", font=M, spacing=17)
text(s, 0.95, 3.50, 5.2, 0.3, "zone-b has 561 kW. The request needs 792 kW.",
     size=13, bold=True, color="red")
text(s, 0.95, 3.78, 5.2, 0.28, "Confident, well formed, and wrong by 231 kW.", size=11.5)
for i, (n, d, col) in enumerate([
        ("Guardrail", "Re-checks every verdict against the same figures the assessor was given. "
                      "Not the model checking itself. Ordinary code.", "grn"),
        ("Retry with correction", "The violation text goes back in: zone-b has 561 kW against a "
                                  "792 kW request. It usually corrects on attempt two.", "cy"),
        ("Fail safe, never open", "Retries exhausted, so it returns infeasible and escalates. A "
                                  "broken assessor stops the line rather than approving.", "amb")]):
    y = 1.90 + i * 1.45
    card(s, 6.8, y, 5.9, 1.28)
    dot(s, 7.05, y + 0.26, 0.34, col)
    text(s, 7.05, y + 0.30, 0.34, 0.28, str(i + 1), size=13, bold=True, color="bg",
         font=H, align=PP_ALIGN.CENTER)
    text(s, 7.55, y + 0.20, 4.9, 0.3, n, size=14, bold=True, color=col)
    text(s, 7.55, y + 0.52, 4.95, 0.7, d, size=11.5, spacing=16)
card(s, 0.6, 4.30, 5.9, 2.0)
text(s, 0.95, 4.50, 5.2, 0.3, "It caught a failure we had not thought of",
     size=14.5, bold=True, color="ink", font=H)
text(s, 0.95, 4.86, 5.2, 1.4,
     "Twice during the build, a gateway bug broke every request in production. Every assessor hit "
     "its retry limit, refused, and the system returned escalated, needs a human.\n\n"
     "It did not crash and it did not approve anything. The fail safe held against a bug that was "
     "not in the assessors at all.", size=12, spacing=17)
notes(s, "This slide answers a rubric question directly. Six of six bad verdicts are blocked, "
         "including one that hid an impossible zone in its proposed alternative.")

# ============================================ 10. layers of enforcement
s = slide()
title(s, "Five layers, each one verified",
      "Every denial below is produced by Google Cloud, not by our application code")
layers = [
    ("Identity", "Seven service accounts and three custom least-privilege roles. No shared "
                 "credentials anywhere in the system.", "cy"),
    ("Network", "Six of the seven services sit inside a VPC with ingress=internal. An "
                "unauthenticated call gets 404, not 403, so the service is not discoverable.", "vio"),
    ("Data", "IAM conditions on resource.name. A cross-domain read is refused by Firestore "
             "before any of our code runs.", "grn"),
    ("Content", "Model Armor screens traffic between agents on the way in and on the way out. "
                "It fails closed.", "amb"),
    ("Audit", "One correlation id threads every model call, guardrail result, screening result "
              "and routing decision into Cloud Logging.", "red"),
]
for i, (n, d, col) in enumerate(layers):
    y = 1.78 + i * 0.92
    card(s, 0.6, y, 12.1, 0.8)
    dot(s, 0.92, y + 0.23, 0.34, col)
    text(s, 1.42, y + 0.14, 2.3, 0.3, n, size=14, bold=True, color="ink")
    text(s, 3.8, y + 0.16, 8.6, 0.52, d, size=12, spacing=17)
text(s, 0.6, 6.52, 12.1, 0.36,
     "Someone who fully owns the public container gets a request queue and a decision log. Not "
     "the facility, and no way to spend the inference budget.", size=12, color="cy")
notes(s, "The last line is the one worth saying out loud. It is what least privilege buys you.")

# ============================================== 11. running in production
s = slide()
title(s, "Running in production", "Deployed, reproducible, and checkable by anyone")
for i, (n, d, col) in enumerate([
        ("Live on Google Cloud", "Seven Cloud Run services, each under its own service account, "
                                 "scale to zero, instance count capped.", "cy"),
        ("Reproducible from scratch", "Every resource has a script. APIs, databases, IAM roles "
                                      "and conditions, VPC, deploy, seed.", "grn"),
        ("Auditable by design", "A whole multi-agent decision can be reconstructed from one log "
                                "filter on the correlation id.", "vio"),
        ("Decision report", "A PDF per decision showing what each assessor was shown beside what "
                            "it concluded.", "amb")]):
    x = 0.6 + (i % 2) * 6.2
    y = 1.78 + (i // 2) * 1.56
    card(s, x, y, 5.9, 1.36)
    dot(s, x + 0.28, y + 0.30, 0.36, col)
    text(s, x + 0.84, y + 0.24, 4.8, 0.3, n, size=14, bold=True, color="ink")
    text(s, x + 0.84, y + 0.56, 4.9, 0.72, d, size=11.5, spacing=16)
card(s, 0.6, 4.95, 12.1, 1.55, "panel2")
text(s, 0.95, 5.12, 5.0, 0.3, "Run these yourself", size=12.5, bold=True, color="ink")
text(s, 0.95, 5.46, 8.6, 0.9,
     "bash infra/iam/03_verify_isolation.sh    14 database isolation checks\n"
     "bash scripts/demo_denial.sh              every layer, live denials\n"
     "python -m scripts.demo_guardrails        6 guardrail assertions, no model calls",
     size=11, color="grn", font=M, spacing=17)
text(s, 9.85, 5.24, 2.6, 0.5, "60 / day", size=27, bold=True, color="cy", font=H)
text(s, 9.85, 5.78, 2.7, 0.5,
     "hard cap on the public demo, about $1.50 a day worst case", size=11, spacing=15)
notes(s, "The video shows all three of these running unedited. That is the proof of action the "
         "rubric asks for.")

# =================================================== 12. GEAP coverage
s = slide()
title(s, "Gemini Enterprise Agent Platform coverage",
      "Six of the seven capabilities are built. One is partial, and we say so.")
for i, (n, d, st, col) in enumerate([
        ("Agent Identity", "Seven service accounts, three custom roles, IAM conditions", "BUILT", "grn"),
        ("Agent Gateway", "Routing table, identity checks, audience-scoped tokens, audit log", "BUILT", "grn"),
        ("Model Armor", "Screens traffic between agents in both directions, fails closed", "BUILT", "grn"),
        ("Agent Registry", "Contract, version, owner, guardrails and data scope per agent", "BUILT", "grn"),
        ("Memory Bank", "Precedent recalled before reconciling, written back afterwards", "BUILT", "grn"),
        ("Agent Observability", "Correlation id through every agent, full reasoning chain kept", "BUILT", "grn"),
        ("Agent Runtime", "Cloud Run with per-agent identity, but delivery is synchronous", "PARTIAL", "amb")]):
    y = 1.76 + i * 0.66
    card(s, 0.6, y, 12.1, 0.56, "panel2" if i == 6 else "panel")
    text(s, 0.9, y + 0.15, 2.8, 0.28, n, size=13, bold=True, color="ink")
    text(s, 3.75, y + 0.16, 7.1, 0.28, d, size=11.5)
    text(s, 10.9, y + 0.16, 1.5, 0.28, st, size=11.5, bold=True, color=col,
         align=PP_ALIGN.RIGHT)
text(s, 0.6, 6.44, 12.1, 0.44,
     "Each capability is built on Google Cloud primitives rather than adopted as a managed "
     "product. The behaviour is real and checkable, but it is our implementation, and this deck "
     "says so.", size=12, color="faint")
notes(s, "Being straight about the partial one buys credibility for the other six.")

# ================================================ 13. limits and next
s = slide()
title(s, "Limits, and what comes next")
card(s, 0.6, 1.72, 5.9, 4.5)
text(s, 0.95, 1.94, 5.2, 0.3, "Next", size=12, color="faint")
for i, (n, d) in enumerate([
        ("A read-only capacity audit",
         "Point it at a DCIM export and report the megawatts a site has paid for and cannot use. "
         "No integration risk, and it earns the right to be in the decision loop later."),
        ("Durable async execution",
         "Results are already durable, but delivery is still synchronous. Cloud Tasks closes the "
         "last gap."),
        ("A solver under the reasoning",
         "Comparing 792 kW against 3,053 kW is arithmetic. The model earns its place on "
         "interpretation and negotiation, not on the maths.")]):
    y = 2.34 + i * 1.28
    text(s, 0.95, y, 5.2, 0.3, n, size=13.5, bold=True, color="cy")
    text(s, 0.95, y + 0.32, 5.2, 0.88, d, size=11.5, spacing=16)
card(s, 6.8, 1.72, 5.9, 4.5, "panel2")
text(s, 7.15, 1.94, 5.2, 0.3, "Stated plainly", size=12, color="faint")
text(s, 7.15, 2.34, 5.2, 0.3, "The facility is simulated.", size=14, bold=True, color="amb")
text(s, 7.15, 2.74, 5.2, 2.3,
     "The agents, the infrastructure, the isolation and every denial shown are real. The data "
     "center is not. It is generated from published engineering figures, physically consistent, "
     "and reproducible from a fixed seed.\n\n"
     "How these decisions actually get made inside an operations team is something we would need "
     "to learn by shadowing one, not by writing a seed file. That is the honest boundary of what "
     "we built.", size=12, spacing=18)
text(s, 7.15, 5.42, 5.2, 0.7,
     "What is real: 21 security checks, 7 deployed services, and every decision auditable.",
     size=12, bold=True, color="grn", spacing=17)
notes(s, "Judges trust a team more when it draws its own boundary before they have to.")

# ========================================================= 14. close
s = slide()
text(s, 0.9, 2.30, 11.5, 1.4,
     "Four correct answers still add up to nothing\nunless something checks them together.",
     size=29, bold=True, color="ink", font=H, spacing=42)
text(s, 0.95, 3.95, 11.4, 0.4,
     "GridMind, capacity planning for high-density data centers", size=16, color="cy")
for i, (k, v) in enumerate([("Live", "gridmind-wuvfpvopoq-uk.a.run.app"),
                            ("Code", "github.com/PRACHIT27/GridMindAI"),
                            ("Track", "The Fortified Enterprise Fleet")]):
    x = 0.9 + i * 4.0
    text(s, x, 4.95, 3.7, 0.28, k, size=11, color="faint")
    text(s, x, 5.22, 3.8, 0.32, v, size=13, bold=True, color="ink")
notes(s, "Close on the one line, not a thank-you slide.")

prs.save(OUT)
print(f"wrote {OUT}  ({OUT.stat().st_size:,} bytes, {len(prs.slides._sldIdLst)} slides)")
